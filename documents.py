"""Plain text out of the file formats an account manager actually has.

Readers are imported lazily so a cold start doesn't pay for four parsers, and a
problem with one format can't stop the app booting.
"""

import re
import unicodedata
from io import BytesIO
from pathlib import Path
from typing import NamedTuple

# Roughly 12k tokens — a long QBR deck or a 20-page report. The model's context
# window is not the binding constraint here; signal dilution is. Twelve useful
# lines buried in 300 pages extract worse than the same twelve on their own.
# Anything genuinely larger wants retrieval, not a bigger prompt.
MAX_CHARS = 50_000


class Extracted(NamedTuple):
    text: str
    total_chars: int
    truncated: bool


def _pdf(blob: bytes) -> str:
    from pypdf import PdfReader

    return "\n".join(page.extract_text() or "" for page in PdfReader(BytesIO(blob)).pages)


def _docx(blob: bytes) -> str:
    import docx

    document = docx.Document(BytesIO(blob))
    parts = [p.text for p in document.paragraphs]
    # Account plans put half their content in tables.
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _xlsx(blob: bytes) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(BytesIO(blob), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"# {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _pptx(blob: bytes) -> str:
    from pptx import Presentation

    parts: list[str] = []
    for number, slide in enumerate(Presentation(BytesIO(blob)).slides, start=1):
        parts.append(f"# Slide {number}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _plain(blob: bytes) -> str:
    return blob.decode("utf-8", errors="replace")


READERS = {
    ".pdf": _pdf,
    ".docx": _docx,
    ".xlsx": _xlsx,
    ".xlsm": _xlsx,
    ".pptx": _pptx,
    ".csv": _plain,
    ".tsv": _plain,
    ".txt": _plain,
    ".md": _plain,
    ".json": _plain,
}

SUPPORTED = ", ".join(sorted(READERS))


# Zero-width, control and surrogate characters: remove outright.
_DELETE = {"Cf", "Cs", "Cc"}
# Private use: PDF exporters map real glyphs here with no ToUnicode table, so
# "(+50% YoY)" extracts as unrenderable boxes. The mapping is font-specific and
# can't be reversed, but these were visible characters — replace them with a
# space so the words either side don't run together.
_BLANK = "Co"


def clean(text: str) -> str:
    """Drop characters no reader can use, and tidy the whitespace they leave."""
    kept = []
    for ch in text:
        if ch in "\n\t":
            kept.append(ch)
            continue
        category = unicodedata.category(ch)
        if category in _DELETE:
            continue
        kept.append(" " if category == _BLANK else ch)

    collapsed = re.sub(r"[ \t]{2,}", " ", "".join(kept))
    return re.sub(r"\n{3,}", "\n\n", collapsed).strip()


def to_text(filename: str, blob: bytes) -> Extracted:
    """Read `blob` as text, capped at MAX_CHARS. Raises ValueError on an unreadable format.

    Truncation is reported rather than applied silently — a caller that quietly
    drops nine tenths of a contract looks like it worked.
    """
    suffix = Path(filename).suffix.lower()
    reader = READERS.get(suffix)
    if reader is None:
        raise ValueError(f"Can't read '{suffix or filename}'. Supported: {SUPPORTED}")

    try:
        text = reader(blob)
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"That {suffix} file could not be read: {e}") from e

    cleaned = clean(text)
    return Extracted(cleaned[:MAX_CHARS], len(cleaned), len(cleaned) > MAX_CHARS)


def _self_check() -> None:
    assert to_text("notes.txt", b"hello").text == "hello"

    # Private use characters from a PDF font with no ToUnicode table.
    assert clean("Current ARR $1.8M \ue081\ue09d50% YoY\ue082") == "Current ARR $1.8M 50% YoY"
    assert clean("word\ue000word") == "word word", "private use glyphs must not join words"
    assert clean("a\u200bb") == "ab", "zero-width characters go"
    assert clean("keep \u2713 \u2022 \u2014 \u00b7") == "keep \u2713 \u2022 \u2014 \u00b7"
    assert clean("line\n\n\n\nnext") == "line\n\nnext"
    assert clean("wide    gap") == "wide gap"
    assert to_text("a.csv", b"x,y\n1,2").text == "x,y\n1,2"

    small = to_text("ok.md", b"a" * 100)
    assert small.truncated is False and small.total_chars == 100

    big = to_text("big.md", b"a" * (MAX_CHARS + 500))
    assert big.text == "a" * MAX_CHARS
    assert big.truncated is True
    assert big.total_chars == MAX_CHARS + 500

    for bad, expect in [("deck.key", "Supported"), ("old.doc", "Supported")]:
        try:
            to_text(bad, b"")
            raise AssertionError(f"{bad} should not be readable")
        except ValueError as e:
            assert expect in str(e), e

    # A real docx, round-tripped through the writer.
    import docx

    document = docx.Document()
    document.add_paragraph("ARR is $1.4M, flat since last year")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "NPS"
    table.rows[0].cells[1].text = "7"
    buffer = BytesIO()
    document.save(buffer)

    text = to_text("plan.docx", buffer.getvalue()).text
    assert "ARR is $1.4M, flat since last year" in text, text
    assert "NPS | 7" in text, text

    # Corrupt bytes must raise cleanly rather than blowing up.
    try:
        to_text("broken.pdf", b"not a pdf at all")
        raise AssertionError("corrupt pdf should raise")
    except ValueError as e:
        assert "could not be read" in str(e), e

    print("documents self-check passed")


if __name__ == "__main__":
    _self_check()
